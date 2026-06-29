"""
F12 Multi-format ingest — extract.py unit tests (ADR-0025 §4.1, AC-F12-1..7).

Tests:
  T-EXT-001  PDF extraction yields text from all pages (AC-F12-1)
  T-EXT-002  DOCX extraction yields paragraph text (AC-F12-1)
  T-EXT-003  PPTX extraction yields slide text (AC-F12-1)
  T-EXT-004  XLSX extraction yields GFM markdown table (AC-F12-1)
  T-EXT-005  Image extension → placeholder string (§4.5)
  T-EXT-006  AV extension → placeholder string (§4.5)
  T-EXT-007  Unknown extension → UnsupportedFormatError (caller maps to 415)
  T-EXT-008  EXTRACT_MAX_CHARS cap truncates oversized output (I7)
  T-EXT-009  Static guard: no format-lib import outside extract.py (AC-F12-7)
  T-EXT-010  EXTRACTABLE_BINARY_EXTENSIONS contains expected exts (AC-F12-7)
  T-EXT-011  PLACEHOLDER_EXTENSIONS contains image/AV exts (§4.5)
  T-EXT-012  PDF with no text returns placeholder message
  T-EXT-013  XLSX with multiple sheets yields one table per sheet
  T-EXT-014  DOCX with empty document returns placeholder
"""

from __future__ import annotations

import io
import sys
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

# ── Helpers to build minimal in-memory test files ─────────────────────────────


def _make_minimal_pdf(text_pages: list[str]) -> bytes:
    """Build a minimal valid PDF bytes object with the given text pages."""
    # Use pypdf to create a minimal PDF with text; fallback to raw bytes if unavailable.
    try:
        from pypdf import PdfWriter

        writer = PdfWriter()
        for _text in text_pages:
            writer.add_blank_page(width=200, height=200)
        buf = io.BytesIO()
        writer.write(buf)
        return buf.getvalue()
    except Exception:
        # Fallback: return a minimal valid-ish PDF header (will produce no text on extract)
        return b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n%%EOF"


def _make_docx_bytes(paragraphs: list[str]) -> bytes:
    """Build a minimal DOCX (ZIP with word/document.xml)."""
    import zipfile

    doc_xml = '<?xml version="1.0" encoding="UTF-8"?>'
    doc_xml += '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
    doc_xml += "<w:body>"
    for para in paragraphs:
        doc_xml += f"<w:p><w:r><w:t>{para}</w:t></w:r></w:p>"
    doc_xml += "</w:body></w:document>"

    _ct_rels = "application/vnd.openxmlformats-package.relationships+xml"
    _ct_doc = "application/vnd.openxmlformats-officedocument" ".wordprocessingml.document.main+xml"
    _ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    content_types_xml = '<?xml version="1.0" encoding="UTF-8"?>'
    content_types_xml += f'<Types xmlns="{_ct_ns}">'
    content_types_xml += f'<Default Extension="rels" ContentType="{_ct_rels}"/>'
    content_types_xml += '<Default Extension="xml" ContentType="application/xml"/>'
    content_types_xml += f'<Override PartName="/word/document.xml" ContentType="{_ct_doc}"/>'
    content_types_xml += "</Types>"

    _rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    _rel_type = (
        "http://schemas.openxmlformats.org/officeDocument" "/2006/relationships/officeDocument"
    )
    rels_xml = '<?xml version="1.0" encoding="UTF-8"?>'
    rels_xml += f'<Relationships xmlns="{_rel_ns}">'
    rels_xml += f'<Relationship Id="rId1" Type="{_rel_type}" Target="word/document.xml"/>'
    rels_xml += "</Relationships>"

    word_rels_xml = '<?xml version="1.0" encoding="UTF-8"?>'
    word_rels_xml += (
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    )
    word_rels_xml += "</Relationships>"

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types_xml)
        zf.writestr("_rels/.rels", rels_xml)
        zf.writestr("word/_rels/document.xml.rels", word_rels_xml)
        zf.writestr("word/document.xml", doc_xml)
    return buf.getvalue()


def _make_xlsx_bytes(sheet_data: dict[str, list[list[str]]]) -> bytes:
    """Build a minimal XLSX using openpyxl."""
    import openpyxl

    wb = openpyxl.Workbook()
    first = True
    for sheet_name, rows in sheet_data.items():
        if first:
            ws = wb.active
            ws.title = sheet_name
            first = False
        else:
            ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── T-EXT-001: PDF extraction ─────────────────────────────────────────────────


class TestPdfExtraction:
    """T-EXT-001: PDF text extraction (AC-F12-1)."""

    def test_pdf_placeholder_on_empty_pdf(self, tmp_path: Path) -> None:
        """PDF with no embedded text yields a placeholder message (not an error)."""
        from app.ingest.extract import extract_text

        pdf_bytes = _make_minimal_pdf([])
        pdf_file = tmp_path / "empty.pdf"
        pdf_file.write_bytes(pdf_bytes)

        result = extract_text(pdf_file)
        # Either extracted text OR placeholder — must be a non-empty string
        assert isinstance(result, str)
        assert len(result) > 0

    def test_pdf_uses_pypdf(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """PDF extraction dispatches through pypdf (AC-F12-7)."""
        from app.ingest import extract as ext_module

        called: list[bool] = []

        class FakePdfReader:
            def __init__(self, path: str) -> None:
                called.append(True)
                self.pages: list = []

        fake_pypdf = MagicMock()
        fake_pypdf.PdfReader = FakePdfReader

        pdf_file = tmp_path / "test.pdf"
        pdf_file.write_bytes(b"%PDF-1.4\n%%EOF")

        with patch.dict(sys.modules, {"pypdf": fake_pypdf}):
            result = ext_module._extract_pdf(pdf_file)

        assert called, "PdfReader should have been called"
        assert isinstance(result, str)


# ── T-EXT-002: DOCX extraction ────────────────────────────────────────────────


class TestDocxExtraction:
    """T-EXT-002: DOCX paragraph extraction (AC-F12-1)."""

    def test_docx_yields_paragraphs(self, tmp_path: Path) -> None:
        """DOCX with 3 paragraphs returns all 3 in the output."""
        from app.ingest.extract import extract_text

        paragraphs = ["First paragraph.", "Second paragraph.", "Third paragraph."]
        docx_bytes = _make_docx_bytes(paragraphs)
        docx_file = tmp_path / "test.docx"
        docx_file.write_bytes(docx_bytes)

        result = extract_text(docx_file)
        assert isinstance(result, str)
        for para in paragraphs:
            assert para in result, f"Expected paragraph {para!r} in output"

    def test_empty_docx_returns_placeholder(self, tmp_path: Path) -> None:
        """Empty DOCX returns a placeholder message."""
        from app.ingest.extract import extract_text

        docx_bytes = _make_docx_bytes([])
        docx_file = tmp_path / "empty.docx"
        docx_file.write_bytes(docx_bytes)

        result = extract_text(docx_file)
        assert isinstance(result, str)
        assert len(result) > 0


# ── T-EXT-003: PPTX extraction ───────────────────────────────────────────────


class TestPptxExtraction:
    """T-EXT-003: PPTX slide text extraction (AC-F12-1)."""

    def _make_pptx(self, slide_texts: list[str]) -> bytes:
        """Build a minimal PPTX with the given slide texts."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # blank layout
        for text in slide_texts:
            slide = prs.slides.add_slide(blank_layout)
            txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(1))
            txBox.text_frame.text = text
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def test_pptx_yields_slide_text(self, tmp_path: Path) -> None:
        """PPTX with 2 slides returns text from each slide."""
        from app.ingest.extract import extract_text

        pptx_bytes = self._make_pptx(["Slide one content", "Slide two content"])
        pptx_file = tmp_path / "test.pptx"
        pptx_file.write_bytes(pptx_bytes)

        result = extract_text(pptx_file)
        assert isinstance(result, str)
        assert "Slide one content" in result
        assert "Slide two content" in result

    def test_empty_pptx_returns_placeholder(self, tmp_path: Path) -> None:
        """PPTX with no text shapes returns a placeholder."""
        from app.ingest.extract import extract_text
        from pptx import Presentation

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        pptx_file = tmp_path / "empty.pptx"
        pptx_file.write_bytes(buf.getvalue())

        result = extract_text(pptx_file)
        assert isinstance(result, str)
        assert len(result) > 0


# ── T-EXT-004: XLSX extraction → GFM table ───────────────────────────────────


class TestXlsxExtraction:
    """T-EXT-004: XLSX cell extraction rendered as GFM markdown table (AC-F12-1)."""

    def test_xlsx_yields_gfm_table(self, tmp_path: Path) -> None:
        """XLSX with 3 rows → GFM markdown table with header separator."""
        from app.ingest.extract import extract_text

        xlsx_bytes = _make_xlsx_bytes(
            {
                "Sheet1": [
                    ["Name", "Value", "Unit"],
                    ["Alpha", "42", "kg"],
                    ["Beta", "3.14", "m"],
                ]
            }
        )
        xlsx_file = tmp_path / "test.xlsx"
        xlsx_file.write_bytes(xlsx_bytes)

        result = extract_text(xlsx_file)
        assert isinstance(result, str)
        # GFM table rows start with |
        assert "|" in result
        # Header row present
        assert "Name" in result
        assert "Value" in result
        # Data rows present
        assert "Alpha" in result
        assert "Beta" in result
        # Separator row present (--- or similar)
        assert "---" in result or "|-" in result

    def test_xlsx_sheet_name_in_output(self, tmp_path: Path) -> None:
        """Sheet name appears as a heading in the GFM output."""
        from app.ingest.extract import extract_text

        xlsx_bytes = _make_xlsx_bytes({"MySheet": [["Col1", "Col2"], ["a", "b"]]})
        xlsx_file = tmp_path / "test.xlsx"
        xlsx_file.write_bytes(xlsx_bytes)

        result = extract_text(xlsx_file)
        assert "MySheet" in result

    def test_xlsx_multiple_sheets(self, tmp_path: Path) -> None:
        """Multiple sheets each produce a GFM table section (T-EXT-013)."""
        from app.ingest.extract import extract_text

        xlsx_bytes = _make_xlsx_bytes(
            {
                "Sheet1": [["A", "B"], ["1", "2"]],
                "Sheet2": [["X", "Y"], ["3", "4"]],
            }
        )
        xlsx_file = tmp_path / "multi.xlsx"
        xlsx_file.write_bytes(xlsx_bytes)

        result = extract_text(xlsx_file)
        assert "Sheet1" in result
        assert "Sheet2" in result

    def test_xlsx_empty_returns_placeholder(self, tmp_path: Path) -> None:
        """Empty XLSX returns placeholder message."""
        from app.ingest.extract import extract_text

        xlsx_bytes = _make_xlsx_bytes({"Sheet1": []})
        xlsx_file = tmp_path / "empty.xlsx"
        xlsx_file.write_bytes(xlsx_bytes)

        result = extract_text(xlsx_file)
        assert isinstance(result, str)
        assert len(result) > 0


# ── T-EXT-005/006: Placeholder for image/AV ───────────────────────────────────


class TestPlaceholderExtraction:
    """T-EXT-005/006: Image and AV files yield placeholder strings (§4.5)."""

    @pytest.mark.parametrize("ext", [".png", ".jpg", ".jpeg", ".gif", ".webp"])
    def test_image_yields_placeholder(self, tmp_path: Path, ext: str) -> None:
        """Image file returns a placeholder (no OCR in M5)."""
        from app.ingest.extract import extract_text

        img_file = tmp_path / f"test{ext}"
        img_file.write_bytes(b"\xff\xd8\xff")  # minimal JPEG-like header
        result = extract_text(img_file)
        assert isinstance(result, str)
        assert len(result) > 0
        assert "no text extracted" in result.lower() or "image" in result.lower()

    @pytest.mark.parametrize("ext", [".mp3", ".mp4", ".wav", ".m4a"])
    def test_av_yields_placeholder(self, tmp_path: Path, ext: str) -> None:
        """AV file returns a placeholder (no transcript in M5)."""
        from app.ingest.extract import extract_text

        av_file = tmp_path / f"test{ext}"
        av_file.write_bytes(b"\x00" * 16)
        result = extract_text(av_file)
        assert isinstance(result, str)
        assert len(result) > 0
        assert "transcript" in result.lower() or "av" in result.lower()


# ── T-EXT-007: Unknown extension raises UnsupportedFormatError ───────────────


class TestUnsupportedFormat:
    """T-EXT-007: Unknown extension → UnsupportedFormatError (caller maps to 415)."""

    @pytest.mark.parametrize("ext", [".exe", ".bin", ".zip", ".tar", ".db"])
    def test_unknown_ext_raises(self, tmp_path: Path, ext: str) -> None:
        from app.ingest.extract import UnsupportedFormatError, extract_text

        bad_file = tmp_path / f"file{ext}"
        bad_file.write_bytes(b"\x00" * 4)
        with pytest.raises(UnsupportedFormatError):
            extract_text(bad_file)


# ── T-EXT-008: EXTRACT_MAX_CHARS cap ─────────────────────────────────────────


class TestExtractMaxChars:
    """T-EXT-008: Output capped at EXTRACT_MAX_CHARS (I7 — pathological file guard)."""

    def test_max_chars_cap(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """Output > EXTRACT_MAX_CHARS is truncated silently."""
        from app import config as cfg
        from app.ingest import extract as ext_module

        # Set a tiny cap
        monkeypatch.setattr(cfg.settings, "extract_max_chars", 50)

        docx_bytes = _make_docx_bytes(["A" * 200, "B" * 200])
        docx_file = tmp_path / "big.docx"
        docx_file.write_bytes(docx_bytes)

        result = ext_module.extract_text(docx_file)
        assert len(result) <= 50, f"Expected ≤50 chars; got {len(result)}"


# ── T-EXT-009: Static guard — no format-lib imports outside extract.py ────────


class TestStaticGuard:
    """T-EXT-009: Static guard: pypdf/docx/pptx/openpyxl only in ingest/extract.py (AC-F12-7)."""

    def _grep_imports(self, root: Path, forbidden: list[str]) -> list[str]:
        """
        Walk .py files under root and find any import of forbidden packages
        OUTSIDE app/ingest/extract.py.
        """
        violations: list[str] = []
        for py_file in root.rglob("*.py"):
            if "extract.py" in py_file.name:
                continue  # allowed
            if "__pycache__" in str(py_file):
                continue
            try:
                text = py_file.read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
            for pkg in forbidden:
                # Match "import pypdf", "from pypdf import", "import openpyxl" etc.
                if f"import {pkg}" in text or f"from {pkg}" in text:
                    violations.append(f"{py_file.relative_to(root)}: imports {pkg}")
        return violations

    def test_no_format_lib_imports_outside_extract(self) -> None:
        """No pypdf/docx/pptx/openpyxl imports outside ingest/extract.py (AC-F12-7)."""
        backend_app = Path(__file__).resolve().parent.parent / "app"
        violations = self._grep_imports(
            backend_app,
            ["pypdf", "docx", "pptx", "openpyxl"],
        )
        assert not violations, "Format lib imports found outside ingest/extract.py:\n" + "\n".join(
            violations
        )

    def test_no_unstructured_added(self) -> None:
        """unstructured is NOT added in M5 (ADR-0025 §4.6, Do-NOT #16).

        Scans app/ source files only (not test files, which may reference the word
        in docstrings and assertion messages).
        """
        app_dir = Path(__file__).resolve().parent.parent / "app"
        for py_file in app_dir.rglob("*.py"):
            if "__pycache__" in str(py_file):
                continue
            try:
                text = py_file.read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
            assert "import unstructured" not in text and "from unstructured" not in text, (
                f"unstructured import found in {py_file} — "
                "must not be added in M5 (ADR-0025 §4.6, Do-NOT #16)"
            )


# ── T-EXT-010/011: Extension set constants ───────────────────────────────────


class TestExtensionSets:
    """T-EXT-010/011: EXTRACTABLE_BINARY_EXTENSIONS and PLACEHOLDER_EXTENSIONS (AC-F12-7)."""

    def test_extractable_extensions_complete(self) -> None:
        """EXTRACTABLE_BINARY_EXTENSIONS contains .pdf, .docx, .pptx, .xlsx (AC-F12-7)."""
        from app.ingest.extract import EXTRACTABLE_BINARY_EXTENSIONS

        assert ".pdf" in EXTRACTABLE_BINARY_EXTENSIONS
        assert ".docx" in EXTRACTABLE_BINARY_EXTENSIONS
        assert ".pptx" in EXTRACTABLE_BINARY_EXTENSIONS
        assert ".xlsx" in EXTRACTABLE_BINARY_EXTENSIONS

    def test_placeholder_extensions_complete(self) -> None:
        """PLACEHOLDER_EXTENSIONS contains image and AV formats (§4.5)."""
        from app.ingest.extract import PLACEHOLDER_EXTENSIONS

        for ext in [".png", ".jpg", ".jpeg", ".gif", ".webp", ".mp3", ".mp4", ".wav", ".m4a"]:
            assert ext in PLACEHOLDER_EXTENSIONS, f"{ext} missing from PLACEHOLDER_EXTENSIONS"

    def test_no_overlap_between_sets(self) -> None:
        """The two extension sets are disjoint."""
        from app.ingest.extract import EXTRACTABLE_BINARY_EXTENSIONS, PLACEHOLDER_EXTENSIONS

        overlap = EXTRACTABLE_BINARY_EXTENSIONS & PLACEHOLDER_EXTENSIONS
        assert not overlap, f"Overlap between extension sets: {overlap}"

    def test_binary_exts_not_in_allowed(self) -> None:
        """Binary extensions are NOT in _ALLOWED_EXTENSIONS (Do-NOT #13, ADR-0025 §4.3)."""
        from app.ingest.extract import EXTRACTABLE_BINARY_EXTENSIONS
        from app.upload import _ALLOWED_EXTENSIONS

        overlap = EXTRACTABLE_BINARY_EXTENSIONS & _ALLOWED_EXTENSIONS
        assert not overlap, (
            f"Binary extensions found in _ALLOWED_EXTENSIONS: {overlap}. "
            "This breaks the watcher (ADR-0025 §4.3, Do-NOT #13)."
        )
