"""
F12 Multi-format text extractor — SOLE home of format-specific libraries (ADR-0025 §4.1).

STATIC GUARD (AC-F12-7 / ADR-0051):
  pypdf, docx, pptx, openpyxl, striprtf, odf, markdownify MUST NOT be imported anywhere
  outside this module.  Any PR that introduces those imports elsewhere is a P0 rejection.

  NOTE (ADR-0051 / R8-1): When PDF_EXTRACTOR=marker, extract.py calls an HTTP
  microservice at MARKER_SERVICE_URL/convert. The marker package itself is NOT
  imported here and does NOT live in the backend container. pypdf is still the
  sole container-side PDF library; Marker is called over HTTP, not imported.

INVARIANT CONTRACT:
  I6: extract_text() is PURE (path in, text out) — one documented exception:
      _extract_pdf_via_marker() makes an HTTP call to the external Marker
      microservice. This is the PM-approved exception for R8-1 (ADR-0051).
      All other paths are inference-free.
  I7: output capped at EXTRACT_MAX_CHARS (config). No loop — single-pass per file.
      The Marker call uses a bounded timeout (MARKER_TIMEOUT_SECONDS).
  I9: Uses well-known pure-Python extractor libs; unstructured deliberately NOT added (§4.5).
      Marker is called over HTTP — not imported — so heavy ML deps stay host-side.
  I5: The companion .extracted.md written by the caller (upload handler) has valid YAML
      frontmatter — this module only returns text; formatting is the caller's responsibility.

Extension dispatch:
  .pdf  → _extract_pdf_via_marker() when PDF_EXTRACTOR=marker (falls back to pypdf on
          any failure); _extract_pdf() (pypdf) when PDF_EXTRACTOR=pypdf (default).
          Images in PDF skipped with WARNING — AC-F12-1.
  .docx → python-docx (paragraphs)
  .pptx → python-pptx (slide text)
  .xlsx → openpyxl (sheets → GFM markdown table)
  .csv  → stdlib csv (NO new dep) → GFM markdown table
  .html → markdownify (pure-python; strips scripts/styles; tables → GFM)
  .mdx  → UTF-8 text with import/export and JSX tag stripping (no dep)
  .rtf  → striprtf (pure-python, rtf_to_text)
  .odt  → odfpy (text:p paragraphs)
  .ods  → odfpy (table cells → GFM markdown table per sheet)
  .odp  → odfpy (draw:page / draw:text-box paragraphs → ## Slide N headers)
  .doc  → NOT implemented. Legacy binary OLE Word format has no trivial pure-python
          extractor. The only well-known pure-python candidates (python-docx2txt,
          antiword, libreoffice) are either system-level binaries or require heavy
          optional C extensions. Decision: defer to a future sprint when a lightweight
          solution is available. .doc files are NOT added to EXTRACTABLE_BINARY_EXTENSIONS.
          (ADR note: P3-c decision, sprint v1.5.)
  image/AV in PLACEHOLDER_EXTENSIONS → one-line placeholder (§4.5; no OCR/transcript)
  anything else → raise UnsupportedFormatError (caller maps to HTTP 415)
"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# ── Extension sets (used by upload.py for type-gating) ───────────────────────
# NOTE: Do NOT add these to upload._ALLOWED_EXTENSIONS — the watcher imports that
# frozenset and must remain format-agnostic (ADR-0025 §4.3, Do-NOT #13).

EXTRACTABLE_BINARY_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        # P3-c additions (v1.5 LLM Wiki parity) [F12]:
        ".csv",  # stdlib csv → GFM table; no new dep
        ".html",  # markdownify → readable markdown
        ".mdx",  # Markdown+JSX → stripped text; no dep
        ".rtf",  # striprtf → plain text
        ".odt",  # odfpy → paragraph text
        ".ods",  # odfpy → GFM table per sheet
        ".odp",  # odfpy → slide text
        # .doc NOT included — see module docstring for ADR note
    }
)
"""Extensions extracted by the format libs (pypdf / python-docx / python-pptx / openpyxl /
csv / markdownify / striprtf / odfpy)."""

PLACEHOLDER_EXTENSIONS: frozenset[str] = frozenset(
    {".png", ".jpg", ".jpeg", ".gif", ".webp", ".mp3", ".mp4", ".wav", ".m4a"}
)
"""Extensions accepted but not transcribed/OCR'd in M5 (§4.5 — placeholder only)."""


class UnsupportedFormatError(ValueError):
    """Raised when the file extension is not in any accepted set (caller maps to HTTP 415)."""


def _extract_max_chars() -> int:
    """Cap on extracted text output (I7 — pathological file guard). From EXTRACT_MAX_CHARS env."""
    try:
        from app.config import settings

        return int(getattr(settings, "extract_max_chars", 2_000_000))
    except Exception:  # noqa: BLE001
        return 2_000_000


def _get_pdf_extractor() -> str:
    """Return the effective PDF extractor backend ('pypdf' or 'marker').

    S1 (ADR-0053 §2.5): read effective value through config_overrides.
    Falls back to settings.pdf_extractor (env baseline) when no override exists.
    """
    try:
        from app.config import settings  # noqa: PLC0415
        from app.config_overrides import effective_str  # noqa: PLC0415

        return (
            str(effective_str("pdf_extractor", settings.pdf_extractor) or "pypdf").lower().strip()
        )
    except Exception:  # noqa: BLE001
        return "pypdf"


def _get_marker_settings() -> tuple[str, float]:
    """Return (effective_marker_service_url, effective_marker_timeout_seconds).

    S2/S3 (ADR-0053 §2.5): read effective values through config_overrides.
    """
    try:
        from app.config import settings  # noqa: PLC0415
        from app.config_overrides import effective_float, effective_str  # noqa: PLC0415

        url = str(
            effective_str("marker_service_url", settings.marker_service_url)
            or "http://host.docker.internal:8555"
        )
        timeout = effective_float("marker_timeout_seconds", settings.marker_timeout_seconds)
        return url, timeout
    except Exception:  # noqa: BLE001
        return "http://host.docker.internal:8555", 120.0


def _extract_pdf_via_marker(path: Path) -> str | None:
    """
    Call the Marker microservice to extract PDF text (ADR-0051 / R8-1).

    POSTs the raw PDF bytes to {MARKER_SERVICE_URL}/convert with a bounded timeout
    (MARKER_TIMEOUT_SECONDS). On success returns the markdown string from the response.
    On ANY failure (connection refused, timeout, non-200, invalid JSON, missing field)
    logs a WARNING and returns None — the caller MUST fall back to pypdf.

    This is the SOLE network call in extract.py and the PM-approved exception to the
    I6 pure-function contract (ADR-0051 §3).

    AC-R8-1-1: correct request shape (multipart 'file' field with PDF bytes + 30-s-class
    timeout); fallback signalled by returning None.
    """
    import httpx  # noqa: PLC0415 — short-lived client; httpx is a backend dependency

    marker_url, timeout = _get_marker_settings()
    convert_url = f"{marker_url.rstrip('/')}/convert"

    try:
        with httpx.Client(timeout=timeout) as client:
            response = client.post(
                convert_url,
                files={"file": (path.name, path.read_bytes(), "application/pdf")},
            )
        if response.status_code != 200:
            logger.warning(
                "extract_pdf_via_marker: non-200 %d from %s for %s — falling back to pypdf",
                response.status_code,
                convert_url,
                path.name,
            )
            return None
        data = response.json()
        markdown = data.get("markdown")
        if not isinstance(markdown, str) or not markdown:
            logger.warning(
                "extract_pdf_via_marker: invalid/empty 'markdown' in response from %s for %s "
                "— falling back to pypdf",
                convert_url,
                path.name,
            )
            return None
        logger.info(
            "extract_pdf_via_marker: extracted %d chars from %s via Marker (%d pages)",
            len(markdown),
            path.name,
            data.get("pages", 0),
        )
        return markdown
    except Exception as exc:  # noqa: BLE001
        logger.warning(
            "extract_pdf_via_marker: call to %s failed for %s: %s — falling back to pypdf",
            convert_url,
            path.name,
            exc,
        )
        return None


def extract_text(file_path: str | Path) -> str:
    """
    Dispatch on the lower-cased file extension and return extracted plain text (ADR-0025 §4.1).

    Output is capped at EXTRACT_MAX_CHARS (I7). For PDFs, dispatches to the Marker
    microservice when PDF_EXTRACTOR=marker; falls back to pypdf unconditionally on any
    failure (ADR-0051). With the default PDF_EXTRACTOR=pypdf the call path is identical
    to pre-v0.8 behaviour (AC-R8-1-2).

    Raises UnsupportedFormatError for extensions not in EXTRACTABLE_BINARY_EXTENSIONS or
    PLACEHOLDER_EXTENSIONS — the caller (upload handler) maps this to HTTP 415.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()
    max_chars = _extract_max_chars()

    if suffix == ".pdf":
        # R8-1 / ADR-0051: dispatch to Marker when configured; unconditional pypdf fallback
        if _get_pdf_extractor() == "marker":
            marker_result = _extract_pdf_via_marker(path)
            text = marker_result if marker_result is not None else _extract_pdf(path)
        else:
            text = _extract_pdf(path)
    elif suffix == ".docx":
        text = _extract_docx(path)
    elif suffix == ".pptx":
        text = _extract_pptx(path)
    elif suffix == ".xlsx":
        text = _extract_xlsx(path)
    # ── P3-c additions (v1.5 LLM Wiki parity) [F12] ──────────────────────────
    elif suffix == ".csv":
        text = _extract_csv(path)
    elif suffix == ".html":
        text = _extract_html(path)
    elif suffix == ".mdx":
        text = _extract_mdx(path)
    elif suffix == ".rtf":
        text = _extract_rtf(path)
    elif suffix == ".odt":
        text = _extract_odt(path)
    elif suffix == ".ods":
        text = _extract_ods(path)
    elif suffix == ".odp":
        text = _extract_odp(path)
    # ─────────────────────────────────────────────────────────────────────────
    elif suffix in PLACEHOLDER_EXTENSIONS:
        text = _placeholder(path)
    else:
        raise UnsupportedFormatError(
            f"Unsupported file format: {suffix!r}. "
            f"Accepted binary formats: {sorted(EXTRACTABLE_BINARY_EXTENSIONS)}; "
            f"placeholder formats: {sorted(PLACEHOLDER_EXTENSIONS)}."
        )

    # I7 output cap — truncate silently (the extracted text may still be useful)
    if len(text) > max_chars:
        logger.warning(
            "extract_text: output truncated from %d to %d chars (EXTRACT_MAX_CHARS) for %s",
            len(text),
            max_chars,
            path.name,
        )
        text = text[:max_chars]

    return text


# ── Per-format extractors ─────────────────────────────────────────────────────
# pypdf / docx / pptx / openpyxl are imported LOCALLY inside each helper so they
# are only loaded when the corresponding format is actually processed. Any import
# error surfaces at extraction time (not at module import), and the static guard
# check only tests for module-level imports.


def _extract_pdf(path: Path) -> str:
    """Extract text from a PDF using pypdf (AC-F12-1: images in PDF skipped)."""
    import pypdf  # noqa: PLC0415 — local import enforces static guard (AC-F12-7)

    reader = pypdf.PdfReader(str(path))
    pages_text: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001
            # Images/forms that cannot be extracted are silently skipped (AC-F12-1)
            logger.warning(
                "extract_pdf: page %d extraction error in %s: %s — skipping", i, path.name, exc
            )
            page_text = ""
        if page_text:
            pages_text.append(page_text)

    if not pages_text:
        logger.warning("extract_pdf: no text extracted from %s (images-only PDF?)", path.name)
        return f"PDF file: no text content extracted from {path.name}."

    return "\n\n".join(pages_text)


def _extract_docx(path: Path) -> str:
    """Extract paragraph text from a DOCX using python-docx."""
    import docx  # noqa: PLC0415 — local import enforces static guard (AC-F12-7)

    doc = docx.Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return f"DOCX file: no text content extracted from {path.name}."
    return "\n\n".join(paragraphs)


def _extract_pptx(path: Path) -> str:
    """Extract slide text from a PPTX using python-pptx (one logical doc)."""
    from pptx import Presentation  # noqa: PLC0415 — local import enforces static guard (AC-F12-7)

    prs = Presentation(str(path))
    slide_texts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
        if parts:
            slide_texts.append(f"## Slide {slide_num}\n\n" + "\n".join(parts))

    if not slide_texts:
        return f"PPTX file: no text content extracted from {path.name}."
    return "\n\n".join(slide_texts)


def _extract_xlsx(path: Path) -> str:
    """
    Extract cell text from an XLSX using openpyxl, rendered as GFM markdown tables.

    One table per sheet. Non-empty rows only.
    """
    import openpyxl  # noqa: PLC0415 — local import enforces static guard (AC-F12-7)

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheet_texts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            # Convert each cell to string; skip rows that are entirely empty
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)

        if not rows:
            continue

        # GFM table: first data row as header, then separator, then remaining rows
        if len(rows) >= 1:
            header = rows[0]
            sep = ["-" * max(len(h), 3) for h in header]
            table_lines = [
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(sep) + " |",
            ]
            for data_row in rows[1:]:
                # Pad shorter rows to header width
                padded = data_row + [""] * max(0, len(header) - len(data_row))
                table_lines.append("| " + " | ".join(padded[: len(header)]) + " |")
            sheet_texts.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(table_lines))

    wb.close()

    if not sheet_texts:
        return f"XLSX file: no cell content extracted from {path.name}."
    return "\n\n".join(sheet_texts)


# ── P3-c extractors (v1.5 LLM Wiki parity) [F12] ────────────────────────────
# All new format libs are imported LOCALLY to enforce the static guard (AC-F12-7).


def _extract_csv(path: Path) -> str:
    """
    Extract text from a CSV file using stdlib csv → GFM markdown table (no new dep).

    First data row is treated as the header.  Empty rows are skipped.
    """
    import csv  # noqa: PLC0415 — stdlib; local import keeps the pattern consistent

    try:
        with path.open(newline="", encoding="utf-8", errors="replace") as fh:
            reader = csv.reader(fh)
            rows: list[list[str]] = [row for row in reader if any(c.strip() for c in row)]
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_csv: failed to read %s: %s", path.name, exc)
        return f"CSV file: could not be read ({path.name})."

    if not rows:
        return f"CSV file: no content extracted from {path.name}."

    header = rows[0]
    sep = ["-" * max(len(h), 3) for h in header]
    table_lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(sep) + " |",
    ]
    for data_row in rows[1:]:
        padded = data_row + [""] * max(0, len(header) - len(data_row))
        table_lines.append("| " + " | ".join(padded[: len(header)]) + " |")

    return "\n".join(table_lines)


def _extract_html(path: Path) -> str:
    """
    Extract readable markdown from an HTML file using markdownify.

    markdownify is pure-python and converts HTML to GFM markdown including
    tables, headings, and links.  <script> and <style> tag content is removed
    entirely via a regex pre-pass before conversion (markdownify's strip= option
    removes the tags but still emits their inner text; the pre-pass prevents that).
    Falls back to a warning message on parse error.
    """
    import re  # noqa: PLC0415 — stdlib; local for pattern clarity

    import markdownify  # noqa: PLC0415 — local import enforces static guard (AC-F12-7)

    try:
        html_text = path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_html: failed to read %s: %s", path.name, exc)
        return f"HTML file: could not be read ({path.name})."

    if not html_text.strip():
        return f"HTML file: no text content extracted from {path.name}."

    # Pre-strip <script>/<style> blocks entirely (tag + content) before markdownify
    # so their text is never emitted. re.DOTALL handles multi-line blocks.
    html_clean = re.sub(
        r"<(script|style)[^>]*>.*?</(script|style)>",
        "",
        html_text,
        flags=re.DOTALL | re.IGNORECASE,
    )

    try:
        md: str = markdownify.markdownify(html_clean, heading_style="ATX")
        md = md.strip()
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_html: markdownify failed for %s: %s", path.name, exc)
        return f"HTML file: text extraction failed for {path.name}."

    if not md:
        return f"HTML file: no text content extracted from {path.name}."
    return md


def _extract_mdx(path: Path) -> str:
    """
    Extract text from an MDX file (Markdown + JSX).

    MDX files are read as UTF-8 text; obvious import/export statements and
    JSX component tags (PascalCase opening and closing) are stripped.  Standard
    Markdown content and JSX literal text content are preserved.  No external
    dependency — pure stdlib regex.
    """
    import re  # noqa: PLC0415 — stdlib; local for pattern clarity

    try:
        content = path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_mdx: failed to read %s: %s", path.name, exc)
        return f"MDX file: could not be read ({path.name})."

    if not content.strip():
        return f"MDX file: no text content extracted from {path.name}."

    # Strip import and export declaration lines
    content = re.sub(r"^(import|export)\s+.*$", "", content, flags=re.MULTILINE)
    # Strip self-closing JSX tags: <ComponentName ... />
    content = re.sub(r"<[A-Z][A-Za-z0-9.]*[^>]*/\s*>", "", content)
    # Strip opening JSX tags: <ComponentName ...>
    content = re.sub(r"<[A-Z][A-Za-z0-9.]*[^>]*>", "", content)
    # Strip closing JSX tags: </ComponentName>
    content = re.sub(r"</[A-Z][A-Za-z0-9.]*>", "", content)
    # Collapse consecutive blank lines to at most two
    content = re.sub(r"\n{3,}", "\n\n", content).strip()

    if not content:
        return f"MDX file: no text content extracted from {path.name}."
    return content


def _extract_rtf(path: Path) -> str:
    """
    Extract plain text from an RTF file using striprtf (pure-python).

    striprtf.rtf_to_text decodes RTF control words and returns a plain
    unicode string.  encoding errors are silently ignored to handle
    malformed/legacy RTF files.
    """
    from striprtf.striprtf import rtf_to_text  # noqa: PLC0415 — local import enforces guard

    try:
        raw = path.read_bytes()
        # RTF files are typically cp1252; let striprtf detect the codepage directive
        rtf_str = raw.decode("utf-8", errors="replace")
        text = rtf_to_text(rtf_str, errors="ignore")
        text = text.strip()
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_rtf: failed for %s: %s", path.name, exc)
        return f"RTF file: text extraction failed for {path.name}."

    if not text:
        return f"RTF file: no text content extracted from {path.name}."
    return text


def _extract_odt(path: Path) -> str:
    """
    Extract paragraph text from an ODT (OpenDocument Text) file using odfpy.

    Collects all text:p elements and joins non-empty ones with double newlines.
    """
    from odf.opendocument import load as odf_load  # noqa: PLC0415 — local import guard
    from odf.teletype import extractText  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    try:
        doc = odf_load(str(path))
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_odt: failed to open %s: %s", path.name, exc)
        return f"ODT file: could not be read ({path.name})."

    paragraphs: list[str] = []
    for para in doc.getElementsByType(P):
        try:
            text = extractText(para).strip()
        except Exception:  # noqa: BLE001, S112 — skip an unreadable paragraph, keep the rest
            continue
        if text:
            paragraphs.append(text)

    if not paragraphs:
        return f"ODT file: no text content extracted from {path.name}."
    return "\n\n".join(paragraphs)


def _extract_ods(path: Path) -> str:
    """
    Extract cell text from an ODS (OpenDocument Spreadsheet) file using odfpy.

    Renders one GFM markdown table per sheet (same pattern as _extract_xlsx).
    """
    from odf import table as odf_table  # noqa: PLC0415 — local import guard
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.teletype import extractText  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    try:
        doc = odf_load(str(path))
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_ods: failed to open %s: %s", path.name, exc)
        return f"ODS file: could not be read ({path.name})."

    sheet_texts: list[str] = []

    for sheet in doc.getElementsByType(odf_table.Table):
        try:
            sheet_name = sheet.getAttribute("name") or "Sheet"
        except Exception:  # noqa: BLE001
            sheet_name = "Sheet"

        rows: list[list[str]] = []
        for row in sheet.getElementsByType(odf_table.TableRow):
            cells: list[str] = []
            for cell in row.getElementsByType(odf_table.TableCell):
                try:
                    paras = cell.getElementsByType(P)
                    cell_text = " ".join(extractText(p) for p in paras).strip()
                except Exception:  # noqa: BLE001
                    cell_text = ""
                cells.append(cell_text)
            if any(c for c in cells):
                rows.append(cells)

        if not rows:
            continue

        header = rows[0]
        sep = ["-" * max(len(h), 3) for h in header]
        table_lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(sep) + " |",
        ]
        for data_row in rows[1:]:
            padded = data_row + [""] * max(0, len(header) - len(data_row))
            table_lines.append("| " + " | ".join(padded[: len(header)]) + " |")
        sheet_texts.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(table_lines))

    if not sheet_texts:
        return f"ODS file: no cell content extracted from {path.name}."
    return "\n\n".join(sheet_texts)


def _extract_odp(path: Path) -> str:
    """
    Extract slide text from an ODP (OpenDocument Presentation) file using odfpy.

    Collects draw:page elements; within each page gathers text:p elements from
    draw:text-box shapes.  Each slide is prefixed with a ## Slide N header.
    """
    from odf import draw as odf_draw  # noqa: PLC0415 — local import guard
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.teletype import extractText  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    try:
        doc = odf_load(str(path))
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_odp: failed to open %s: %s", path.name, exc)
        return f"ODP file: could not be read ({path.name})."

    slide_texts: list[str] = []

    for slide_num, page in enumerate(doc.getElementsByType(odf_draw.Page), 1):
        parts: list[str] = []
        for tb in page.getElementsByType(odf_draw.TextBox):
            for para in tb.getElementsByType(P):
                try:
                    text = extractText(para).strip()
                except (
                    Exception
                ):  # noqa: BLE001, S112 — skip an unreadable paragraph, keep the rest
                    continue
                if text:
                    parts.append(text)
        if parts:
            slide_texts.append(f"## Slide {slide_num}\n\n" + "\n".join(parts))

    if not slide_texts:
        return f"ODP file: no text content extracted from {path.name}."
    return "\n\n".join(slide_texts)


def _placeholder(path: Path) -> str:
    """
    Return a one-line placeholder for image/AV files (§4.5 — no OCR/transcript in M5).

    Accepted (not 415) but produces a traceable stub; transcription deferred to M6.
    """
    suffix = path.suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        return (
            f"Image file: no text extracted from {path.name}. "
            "Transcription/OCR is out of scope in this release."
        )
    else:
        return (
            f"AV file: transcript not available for {path.name} in this release. "
            "Audio/video transcription is out of scope in this release."
        )
